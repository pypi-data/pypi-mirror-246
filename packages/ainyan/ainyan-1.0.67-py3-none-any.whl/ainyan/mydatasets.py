import re

import datasets
import docx
from simplify_docx import simplify
from pptx import Presentation
from datasets import Dataset, load_dataset, concatenate_datasets
from os.path import splitext
import os
import pandas as pd
import PyPDF2
import argparse

# global
idx = -1

def dataset_from_items(formatted_list):
    filtered_list = [item for item in formatted_list if item['text'] != '' or len(item["text"]) > 3]
    data = {key: [item[key] for item in filtered_list] for key in filtered_list[0]}
    dataset = Dataset.from_dict(data)
    return dataset


def dataset_from_docx(mydoc):
    # read in doc file
    my_doc = docx.Document(mydoc)

    try:
        # coerce to JSON using the standard options
        my_doc_as_json = simplify(my_doc, {"remove-leading-white-space": False})

        # extract body of document
        json_list = my_doc_as_json['VALUE'][0]['VALUE']

        # format to json dataset
        formatted_list = [
            {'url': mydoc, 'index': idx, 'text': item['VALUE'][0]['VALUE'] if item['TYPE'] == 'paragraph' else ""} for
            idx, item in enumerate(json_list)]

        return dataset_from_items(formatted_list)
    except:
        print("Cannot read: " + mydoc)
        return None


def clean_string(input_string):
    # Replace end-of-line characters and tabs with "."
    input_string = re.sub(r'\n|\r|\t', ' ', input_string)
    return input_string

    # Remove all non-alphanumeric characters except ".", space, and Japanese characters
    # cleaned_string = re.sub(r'[^\w. ぁ-んァ-ン一-龯ー ]', '', input_string)
    #
    # return cleaned_string


def split_string(input_str):
    # Regular expression to match English sentences
    english_pattern = re.compile(r'[a-zA-Z.,!? ]+')

    # Regular expression to match Japanese sentences
    japanese_pattern = re.compile(r'[\u3000-\u303F\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FFF、。]+')

    # Find all English and Japanese sentences in the input string
    english_sentences = english_pattern.findall(input_str)
    japanese_sentences = japanese_pattern.findall(input_str)

    return english_sentences + japanese_sentences


def dataset_from_pptx(pptx_file):
    global idx
    # Load the PowerPoint presentation
    presentation = Presentation(pptx_file)

    # Initialize an empty string to store the extracted text
    formatted_list = []

    # Iterate through slides and extract text
    # idx = -1
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                # texts = split_string(text)
                texts = [text]
                for text_ in texts:
                    # print("BEFORE:" + text_)
                    text_ = clean_string(text_)
                    # print("AFTER :" + text_)
                    if len(text_) >= 50:
                        idx = idx + 1
                        print(">>>>>>:"+text_)
                        formatted_list.append({"url": pptx_file, "index": idx, "text": text_})

    return dataset_from_items(formatted_list)


def dataset_from_xlsx(file):
    # Read the Excel file into a dictionary of DataFrames, where keys are sheet names
    xls = pd.read_excel(file, sheet_name=None)
    # df = pd.read_excel(excel_file)

    formatted_list = []

    # idx = -1
    global idx
    for sheet_name, df in xls.items():
        # sheet_text = f"Sheet Name: {sheet_name}\n"  # Include the sheet name
        for column in df.columns:
            for value in df[column]:
                if isinstance(value, str):
                    idx += 1
                    formatted_list.append({"url": file + " " + sheet_name, "index": idx, "text": value + " "})

    return dataset_from_items(formatted_list)


def dataset_from_pdf(file):
    formatted_list = []
    with open(file, 'rb') as pdf_file:
        # Create a PDF reader object
        pdf_reader = PyPDF2.PdfReader(pdf_file)

        # Iterate through each page of the PDF
        # idx = -1
        global idx
        for page_num in range(len(pdf_reader.pages)):
            # Extract the text from the page
            page = pdf_reader.pages[page_num]
            idx += 1
            formatted_list.append({"url": file, "index": idx, "text": page.extract_text()})

    return dataset_from_items(formatted_list)


# private mapping
__mapping = {
    ".doc": dataset_from_docx,  # this actually does not work. doc file can not be loaded directly
    ".docx": dataset_from_docx,
    ".ppt": dataset_from_pptx,
    ".pptx": dataset_from_pptx,
    ".xls": dataset_from_xlsx,
    ".xlsx": dataset_from_xlsx,
    ".pdf": dataset_from_pdf,
}


def dataset_from_file(path):
    root, ext = splitext(path)
    if ext in __mapping:
        target_mapping = __mapping[ext]
        return target_mapping(path)
    else:
        print("No mapping for:" + path)
        return None


def dataset_from_folder(directory_path, include=[], exclude=[]):
    result_list = []

    if not os.path.isdir(directory_path):
        print(directory_path + ": does not exist")
    else:
        for root, dirs, files in os.walk(directory_path):
            for file in files:
                file_path = os.path.join(root, file)
                _, ext = splitext(file_path)
                ext = ext[1:] # remove .
                should_include = (not include and not exclude) or (ext in include) or (exclude and ext not in exclude)
                if should_include:
                    print("processing:" + file_path + "[" + ext + "]")
                    result = dataset_from_file(file_path)
                    if result is not None:
                        result_list.append(result)
                else:
                    print("exclude:" + file_path + "[" + ext + "]")

            for dir_name in dirs:
                subdir_path = os.path.join(root, dir_name)
                result_list.append(dataset_from_folder(subdir_path, include, exclude))

    merged_dataset = concatenate_datasets(result_list)
    return merged_dataset


def dataset_preprocess(source_folder, dataset_save_path, include, exclude):
    """
    Parse all documents in source_folder and convert them to a usable format for training
    :param source_folder: top source folder with docx, xlsx, pptx, .. files
    :param dataset_save_path: where to save the dataset, locally or on S3
    :return: nothing
    """
    dataset = dataset_from_folder(source_folder, include, exclude)
    dataset.save_to_disk(dataset_save_path)


def main():
    parser = argparse.ArgumentParser(description='Two modes: Prepare and upload dataset to s3 or preview the content of a dataset')
    parser.add_argument('--source', metavar='path', required=False,
                        help='the path to the source folder to index files from.')
    parser.add_argument('--dspath', required=True,
                        help='if --source is set, this is the path to save the dataset (either local or s3://), If --source is not set and --preview is then the path to load the dataset from.')
    parser.add_argument('-i', '--include', help="When preparing the dataset, files to include.", nargs='+', default=[])
    parser.add_argument('-e', '--exclude', help="When preparing the dataset, files to exclude", nargs='+', default=[])
    parser.add_argument('--preview', help="Preview dataset switch.", required=False,  action='store_true')
    # only required when
    parser.add_argument('--profile', required=False, help='aws profile to use for upload')
    args = parser.parse_args()

    if args.profile is not None:
        os.environ['AWS_PROFILE'] = args.profile

    if args.source is not None:
        dataset_preprocess(args.source, args.dspath, args.include, args.exclude)

    if args.preview is not None and args.preview is not False:
        ds = datasets.load_from_disk(args.dspath)
        print("{:<8} {:<15} {:<10}".format('Index', 'Url', 'Text'))
        for v in ds.to_iterable_dataset():
            print("{:<8} {:<15} {:<10}".format(v["url"], v["index"], v["text"]))


if __name__ == '__main__':
    main()
